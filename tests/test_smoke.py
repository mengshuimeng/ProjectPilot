from __future__ import annotations

from pathlib import Path

from app.extractor import build_profile
from app.generator import generate_output
from app.llm_client import LLMClient
from app.parser import load_documents, parse_document
from app.pipeline import run_all, run_extract, run_generate
from app.retriever import retrieve_evidence
from app.tool_registry import get_tool_status, search_local_files
from app.verifier import verify_profile


PROJECT_ROOT = Path(__file__).resolve().parent.parent


def _write_fixture_docs(tmp_path: Path) -> Path:
    raw_dir = tmp_path / "raw"
    raw_dir.mkdir()
    (raw_dir / "anchor_project.md").write_text(
        "\n".join(
            [
                "# 通用测试项目",
                "项目名称：通用测试项目",
                "项目背景：该项目面向课程项目材料整理场景。",
                "痛点：项目资料分散，人工整理 README 和答辩稿效率低。",
                "目标用户：学生、教师和项目负责人。",
                "技术路线：Python、Streamlit、LLM、证据检索、自动校验。",
                "系统架构：解析、抽取、检索、生成、校验和重试修复。",
                "创新点：anchor 优先、supporting 补充、verify feedback loop。",
                "交付成果：intro、innovation、defense、readme。",
            ]
        ),
        encoding="utf-8",
    )
    (raw_dir / "supporting_notes.txt").write_text(
        "使用说明：通过 CLI 和 Streamlit 上传材料并生成答辩内容。",
        encoding="utf-8",
    )
    return raw_dir


def test_parser_reads_md_and_pdf(tmp_path: Path) -> None:
    import fitz

    md_path = tmp_path / "project_notes.md"
    md_path.write_text("# 测试项目\n项目背景：用于测试 Markdown 解析。", encoding="utf-8")
    pdf_path = tmp_path / "main.pdf"
    pdf = fitz.open()
    page = pdf.new_page()
    page.insert_text((72, 72), "项目名称：PDF 测试项目")
    pdf.save(pdf_path)
    pdf.close()

    md_doc = parse_document(md_path, role="anchor")
    pdf_doc = parse_document(pdf_path, role="supporting")
    assert md_doc["char_count"] > 0
    assert pdf_doc["suffix"] == ".pdf"
    assert pdf_doc["char_count"] > 0
    assert md_doc["parse_status"] == "parsed"


def test_parser_reads_docx_and_pptx(tmp_path: Path) -> None:
    import docx
    from pptx import Presentation

    docx_path = tmp_path / "main.docx"
    doc = docx.Document()
    doc.add_paragraph("项目名称：通用测试项目")
    doc.add_paragraph("项目背景：用于测试 DOCX 解析。")
    doc.save(docx_path)

    pptx_path = tmp_path / "slides.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "通用测试项目 PPT"
    prs.save(pptx_path)

    docx_doc = parse_document(docx_path, role="anchor")
    pptx_doc = parse_document(pptx_path, role="supporting")
    assert docx_doc["parse_status"] == "parsed"
    assert pptx_doc["parse_status"] == "parsed"
    assert docx_doc["role"] == "anchor"
    assert "通用测试项目" in docx_doc["text"]


def test_extractor_outputs_generic_project_name(tmp_path: Path) -> None:
    raw_dir = _write_fixture_docs(tmp_path)
    docs = load_documents(raw_dir, anchor_name="anchor_project.md")
    profile = build_profile(docs)
    assert profile["project_name"] == "通用测试项目"
    assert profile["anchor_document"] == "anchor_project.md"
    assert "source_roles" in profile


def test_retriever_returns_non_empty_evidence(tmp_path: Path) -> None:
    raw_dir = _write_fixture_docs(tmp_path)
    docs = load_documents(raw_dir, anchor_name="anchor_project.md")
    profile = build_profile(docs)
    evidence = retrieve_evidence("defense", profile, docs)
    assert evidence["task"] == "defense"
    assert evidence["chunks"]
    assert any(chunk["role"] == "anchor" for chunk in evidence["chunks"])


def test_generator_fallback_outputs_string(tmp_path: Path) -> None:
    raw_dir = _write_fixture_docs(tmp_path)
    docs = load_documents(raw_dir, anchor_name="anchor_project.md")
    profile = build_profile(docs)
    content = generate_output(profile, "intro")
    assert isinstance(content, str)
    assert content
    readme = generate_output(profile, "readme")
    assert "# " in readme
    defense = generate_output(profile, "defense")
    assert "这次 V2 升级的目标" not in defense
    assert "ProjectPilot V2 采用了通用" not in defense
    assert "通用测试项目" in defense


def test_llm_client_without_key_does_not_crash(monkeypatch) -> None:
    monkeypatch.delenv("PROJECTPILOT_API_KEY", raising=False)
    monkeypatch.setenv("PROJECTPILOT_LLM_ENABLED", "true")
    result = LLMClient().generate_text("system", "user")
    assert result["ok"] is False
    assert result["fallback"] is True


def test_verifier_returns_passed_key(tmp_path: Path) -> None:
    raw_dir = _write_fixture_docs(tmp_path)
    docs = load_documents(raw_dir, anchor_name="anchor_project.md")
    profile = build_profile(docs)
    report = verify_profile(profile, docs)
    assert isinstance(report, dict)
    assert "passed" in report


def test_runall_outputs_required_files(monkeypatch) -> None:
    monkeypatch.setenv("PROJECTPILOT_LLM_ENABLED", "false")
    run_all(PROJECT_ROOT)
    assert (PROJECT_ROOT / "outputs" / "intro.md").exists()
    assert (PROJECT_ROOT / "outputs" / "innovation.md").exists()
    assert (PROJECT_ROOT / "outputs" / "defense.md").exists()
    assert (PROJECT_ROOT / "outputs" / "readme.md").exists()


def test_generate_reextracts_when_uploaded_material_changes(tmp_path: Path, monkeypatch) -> None:
    monkeypatch.setenv("PROJECTPILOT_LLM_ENABLED", "false")
    project_root = tmp_path
    raw_dir = project_root / "data" / "raw"
    processed_dir = project_root / "data" / "processed"
    raw_dir.mkdir(parents=True)
    processed_dir.mkdir(parents=True)
    (project_root / "outputs").mkdir()

    old_doc = raw_dir / "old_project.md"
    old_doc.write_text("项目名称：旧测试项目\n项目背景：旧项目背景。", encoding="utf-8")
    (processed_dir / "input_manifest.json").write_text(
        '{"anchor_document": "old_project.md", "supporting_documents": []}',
        encoding="utf-8",
    )
    old_profile = run_extract(project_root)
    assert old_profile["project_name"] == "旧测试项目"

    old_doc.unlink()
    new_doc = raw_dir / "new_project.md"
    new_doc.write_text("项目名称：新测试项目\n项目背景：新项目背景。技术路线：Python。", encoding="utf-8")
    (processed_dir / "input_manifest.json").write_text(
        '{"anchor_document": "new_project.md", "supporting_documents": []}',
        encoding="utf-8",
    )

    run_generate(project_root, "intro")
    refreshed_profile = (processed_dir / "profile.json").read_text(encoding="utf-8")
    assert "新测试项目" in refreshed_profile
    assert "old_project.md" not in refreshed_profile


def test_session_workspace_is_isolated(tmp_path: Path, monkeypatch) -> None:
    monkeypatch.setenv("PROJECTPILOT_LLM_ENABLED", "false")
    session_root = tmp_path / "data" / "sessions" / "session-a"
    (session_root / "raw").mkdir(parents=True)
    (session_root / "processed").mkdir()
    (session_root / "outputs").mkdir()
    (session_root / "session.json").write_text('{"session_id": "session-a"}', encoding="utf-8")
    (session_root / "raw" / "anchor.md").write_text(
        "项目名称：会话隔离测试项目\n项目背景：用于验证 session 工作区。",
        encoding="utf-8",
    )
    (session_root / "processed" / "input_manifest.json").write_text(
        '{"anchor_document": "anchor.md", "supporting_documents": []}',
        encoding="utf-8",
    )

    profile = run_extract(session_root)
    result = run_generate(session_root, "intro")
    assert profile["session_id"] == "session-a"
    assert (session_root / "processed" / "profile.json").exists()
    assert (session_root / "outputs" / "intro.md").exists()
    assert "session-a" in (session_root / "outputs" / "intro_meta.json").read_text(encoding="utf-8")
    assert str(session_root) in result["output_path"]


def test_project_name_does_not_include_author_or_abstract() -> None:
    docs = [
        {
            "name": "anchor.md",
            "role": "anchor",
            "text": "基于多模态感知与智能决策的柔性质检分拣机器人系统 张三；李四 摘要 本项目面向智能制造。",
            "char_count": 80,
            "parse_status": "parsed",
            "parse_warning": "",
            "suffix": ".md",
        }
    ]
    profile = build_profile(docs)
    assert profile["project_name"] == "基于多模态感知与智能决策的柔性质检分拣机器人系统"
    assert "摘要" not in profile["project_name"]
    assert "张三" not in profile["project_name"]


def test_verifier_reports_claim_evidence_alignment(tmp_path: Path) -> None:
    raw_dir = _write_fixture_docs(tmp_path)
    docs = load_documents(raw_dir, anchor_name="anchor_project.md")
    profile = build_profile(docs)
    evidence = retrieve_evidence("defense", profile, docs)
    report = verify_profile(
        profile,
        docs,
        output_text="通用测试项目采用 Python 和 Streamlit，实现材料抽取、证据检索和自动校验。",
        evidence=evidence,
        used_sources=[chunk["source"] for chunk in evidence["chunks"]],
        used_roles=[chunk["role"] for chunk in evidence["chunks"]],
        task="defense",
    )
    assert "claim_evidence_alignment" in report
    assert report["claim_evidence_alignment"]["checked_count"] >= 1


def test_local_tool_registry_lists_and_searches_files(tmp_path: Path) -> None:
    (tmp_path / "demo_anchor.md").write_text("项目名称：工具测试项目", encoding="utf-8")
    status = get_tool_status()
    results = search_local_files(tmp_path, query="anchor", suffixes={".md"})
    assert status["tools"]
    assert status["mcp_server_connected"] is False
    assert results and results[0]["name"] == "demo_anchor.md"
