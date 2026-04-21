from __future__ import annotations

import os
import io
from pathlib import Path
from typing import Any, Literal

from app.tool_registry import convert_office_document

try:
    import fitz
except ImportError:  # pragma: no cover - exercised only when dependency is absent
    fitz = None

try:
    import docx
except ImportError:  # pragma: no cover
    docx = None

try:
    from pptx import Presentation
except ImportError:  # pragma: no cover
    Presentation = None

try:
    import pytesseract
except ImportError:  # pragma: no cover
    pytesseract = None

try:
    from PIL import Image
except ImportError:  # pragma: no cover
    Image = None


DocumentRole = Literal["anchor", "supporting"]

TEXT_EXTS = {".md", ".txt"}
STABLE_EXTS = {".txt", ".md", ".pdf", ".docx", ".pptx"}
COMPAT_EXTS = {".doc", ".ppt"}
SUPPORTED_EXTS = STABLE_EXTS | COMPAT_EXTS
OCR_ENABLED = os.getenv("PROJECTPILOT_OCR_ENABLED", "").strip().lower() in {"1", "true", "yes", "on"}
OCR_READY = bool(pytesseract is not None and Image is not None and fitz is not None)


def get_ocr_status() -> dict[str, Any]:
    return {
        "enabled": OCR_ENABLED,
        "available": OCR_READY,
        "note": "OCR fallback 已预留；需安装 pytesseract / Pillow / PyMuPDF 并配置 Tesseract。"
        if OCR_READY
        else "OCR fallback 已预留，但当前环境依赖未就绪。",
    }


def read_text_file(path: Path) -> tuple[str, str]:
    for encoding in ("utf-8", "utf-8-sig", "gb18030", "gbk", "latin-1"):
        try:
            return path.read_text(encoding=encoding), ""
        except UnicodeDecodeError:
            continue
        except Exception as exc:
            return "", f"文本读取失败：{exc}"
    return "", "文本编码无法识别。"


def read_pdf_file(path: Path) -> tuple[str, str]:
    if not path.exists() or path.stat().st_size == 0:
        return "", "文件为空。"
    if fitz is None:
        return "", "PyMuPDF 未安装，无法解析 PDF。请运行 pip install -r requirements.txt。"
    try:
        parts: list[str] = []
        with fitz.open(path) as doc:
            for page in doc:
                parts.append(page.get_text("text"))
        return "\n".join(parts).strip(), ""
    except Exception as exc:
        return "", f"PDF 解析失败：{exc}"


def ocr_pdf_file(path: Path, max_pages: int = 6) -> tuple[str, str]:
    if not OCR_ENABLED:
        return "", "OCR fallback 未启用。"
    if not OCR_READY:
        return "", "OCR 依赖未安装或缺少 PyMuPDF / Pillow / pytesseract。"
    try:
        parts: list[str] = []
        with fitz.open(path) as doc:
            for index, page in enumerate(doc):
                if index >= max_pages:
                    break
                pixmap = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
                image = Image.open(io.BytesIO(pixmap.tobytes("png")))
                text = pytesseract.image_to_string(image, lang="chi_sim+eng")
                if text.strip():
                    parts.append(text.strip())
        if parts:
            return "\n".join(parts).strip(), "PDF 通过 OCR fallback 解析。"
        return "", "OCR 未提取到有效文本。"
    except Exception as exc:
        return "", f"OCR 解析失败：{exc}"


def read_docx_file(path: Path) -> tuple[str, str]:
    if docx is None:
        return "", "python-docx 未安装，无法解析 DOCX。"
    try:
        document = docx.Document(str(path))
        paragraphs = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
        table_texts: list[str] = []
        for table in document.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    table_texts.append(" | ".join(cells))
        return "\n".join(paragraphs + table_texts).strip(), ""
    except Exception as exc:
        return "", f"DOCX 解析失败：{exc}"


def read_pptx_file(path: Path) -> tuple[str, str]:
    if Presentation is None:
        return "", "python-pptx 未安装，无法解析 PPTX。"
    try:
        presentation = Presentation(str(path))
        parts: list[str] = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            slide_parts: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = str(shape.text).strip()
                    if text:
                        slide_parts.append(text)
            if slide_parts:
                parts.append(f"[Slide {slide_index}]\n" + "\n".join(slide_parts))
        return "\n\n".join(parts).strip(), ""
    except Exception as exc:
        return "", f"PPTX 解析失败：{exc}"


def _convert_with_libreoffice(path: Path, target_ext: str) -> tuple[Path | None, str]:
    result = convert_office_document(path, target_ext)
    if not result.get("ok"):
        return None, str(result.get("warning", "旧版 Office 转换失败。"))
    return Path(str(result["converted_path"])), ""


def read_legacy_office_file(path: Path) -> tuple[str, str]:
    if path.suffix.lower() == ".doc":
        converted, warning = _convert_with_libreoffice(path, ".docx")
        if not converted:
            return "", warning
        text, parse_warning = read_docx_file(converted)
        return text, parse_warning or "旧版 DOC 已通过 LibreOffice 转换解析。"
    if path.suffix.lower() == ".ppt":
        converted, warning = _convert_with_libreoffice(path, ".pptx")
        if not converted:
            return "", warning
        text, parse_warning = read_pptx_file(converted)
        return text, parse_warning or "旧版 PPT 已通过 LibreOffice 转换解析。"
    return "", "不支持的旧版 Office 文件。"


def parse_document(path: Path, role: DocumentRole = "supporting") -> dict[str, Any]:
    suffix = path.suffix.lower()
    parse_warning = ""

    if suffix in TEXT_EXTS:
        text, parse_warning = read_text_file(path)
    elif suffix == ".pdf":
        text, parse_warning = read_pdf_file(path)
        if not text.strip():
            ocr_text, ocr_warning = ocr_pdf_file(path)
            if ocr_text.strip():
                text = ocr_text
                parse_warning = ocr_warning
            elif ocr_warning and not parse_warning:
                parse_warning = ocr_warning
    elif suffix == ".docx":
        text, parse_warning = read_docx_file(path)
    elif suffix == ".pptx":
        text, parse_warning = read_pptx_file(path)
    elif suffix in COMPAT_EXTS:
        text, parse_warning = read_legacy_office_file(path)
    else:
        text = ""
        parse_warning = f"不支持的文件类型：{suffix or '无扩展名'}。"

    if text.strip():
        parse_status = "parsed"
    elif suffix not in SUPPORTED_EXTS:
        parse_status = "unsupported"
    elif parse_warning:
        parse_status = "parse_warning"
    else:
        parse_status = "empty"

    return {
        "name": path.name,
        "path": str(path),
        "suffix": suffix,
        "text": text,
        "char_count": len(text),
        "parse_status": parse_status,
        "parse_warning": parse_warning,
        "role": role,
    }


def _anchor_score(path: Path) -> tuple[int, str]:
    name = path.name.lower()
    suffix = path.suffix.lower()
    score = 0
    high_priority_tokens = ["anchor", "main", "primary", "application", "proposal", "方案", "申报", "说明书", "总体"]
    medium_priority_tokens = ["thesis", "paper", "论文", "设计文档", "开发文档"]
    for token in high_priority_tokens:
        if token in name:
            score += 16
    for token in medium_priority_tokens:
        if token in name:
            score += 7
    if suffix in {".pdf", ".docx", ".pptx"}:
        score += 3
    if suffix in TEXT_EXTS:
        score += 1
    try:
        score += min(path.stat().st_size // 100_000, 5)
    except OSError:
        pass
    return score, path.name


def infer_anchor_name(paths: list[Path], anchor_name: str | None = None) -> str:
    if not paths:
        return ""

    explicit = anchor_name or os.getenv("PROJECTPILOT_ANCHOR_DOCUMENT", "").strip()
    if explicit:
        for path in paths:
            if path.name == explicit or path.stem == explicit:
                return path.name

    anchor_dir_paths = [path for path in paths if path.parent.name.lower() == "anchor"]
    if anchor_dir_paths:
        return sorted(anchor_dir_paths)[0].name

    supported_paths = [path for path in paths if path.suffix.lower() in SUPPORTED_EXTS]
    ranked = sorted(supported_paths or paths, key=_anchor_score, reverse=True)
    return ranked[0].name


def load_documents(raw_dir: Path, anchor_name: str | None = None) -> list[dict[str, Any]]:
    if not raw_dir.exists():
        return []

    paths = [
        path
        for path in sorted(raw_dir.rglob("*"))
        if path.is_file() and not path.name.startswith(".")
    ]
    selected_anchor = infer_anchor_name(paths, anchor_name)

    docs: list[dict[str, Any]] = []
    for path in paths:
        role: DocumentRole = "anchor" if path.name == selected_anchor else "supporting"
        docs.append(parse_document(path, role=role))
    return docs
